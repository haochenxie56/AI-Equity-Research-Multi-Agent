"""Thesis Ingestion — LLM extraction layer.

Two LLM entry points:

* :func:`preview_article` — one lightweight call returning the distinct,
  independently-falsifiable arguments an article makes, so the user can pick
  which to extract as separate cards.
* :func:`extract_card` — one call per chosen argument, returning a structured
  (but **not yet validated**) thesis card.

Design rules honoured here:

* All prompts are English; prose fields are emitted bilingually (``*_en`` /
  ``*_zh``) by the model per the new-function convention.
* The LLM only *interprets* the author's text — it never computes a number.
  Every structural / invariant field (IDs, schema versions,
  ``current_evidence_status='unknown'``, ``evidence_refs=[]``,
  ``extraction_meta``) is set **deterministically by this module**, never trusted
  from the model. Validation is delegated to ``validator.py``.
* The only I/O is reading the uploaded document bytes. No network beyond the
  Anthropic client the caller passes in. No import of any scoring / ranking /
  snapshot / anchor module.
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
import os
import pathlib
import re
from datetime import datetime

# Document parsers — now hard dependencies (see requirements.txt). python-docx
# imports as the top-level module ``docx``; python-pptx as ``pptx``.
import docx
import pdfplumber
from dotenv import load_dotenv
from pptx import Presentation

from .schema import (
    PROMPT_VERSION,
    SCHEMA_VERSION_SCENARIO,
    SCHEMA_VERSION_THESIS,
    empty_coi,
)
from .store import card_id_from_hash

_log = logging.getLogger("thesis_ingestion.extractor")

# Model id consistent with the rest of the app's LLM calls (lib/llm_orchestrator).
_MODEL = "claude-sonnet-4-6"


class ExtractionError(Exception):
    """Raised when an LLM extraction response cannot be parsed into JSON."""


class UnsupportedFormatError(Exception):
    """Raised when an uploaded document is in a format we cannot read."""


# ── JSON parsing (pattern copied from lib/llm_orchestrator.py) ────────────────
def _strip_fences(text: str) -> str:
    text = re.sub(r"```(?:json)?\s*", "", text)
    text = re.sub(r"```", "", text)
    return text.strip()


def _repair_json_quotes(text: str) -> str:
    """
    Replace unescaped ASCII double-quotes inside JSON string values
    with Chinese quotation marks to make the JSON parseable.

    Strategy: find patterns where a double-quote appears after a
    non-backslash character that is not a structural JSON character
    (i.e. not after : [ , { or whitespace at the start of a value).
    This is a best-effort heuristic — it handles the most common case
    where the LLM quotes a term like "CXMT" inside a string value.

    Use a state-machine approach: walk the string character by character,
    track whether we are inside a string, and replace unescaped quotes
    that appear mid-string (not at string boundaries).
    """
    result = []
    in_string = False
    i = 0
    while i < len(text):
        ch = text[i]
        if ch == '\\' and in_string:
            # Escaped character — pass through both chars unchanged
            result.append(ch)
            i += 1
            if i < len(text):
                result.append(text[i])
            i += 1
            continue
        if ch == '"':
            if not in_string:
                # Opening quote of a JSON string
                in_string = True
                result.append(ch)
            else:
                # Could be closing quote or unescaped mid-string quote.
                # Peek ahead: if next non-whitespace char is a structural
                # JSON character (: , } ] newline) this is a closing quote.
                j = i + 1
                while j < len(text) and text[j] in ' \t\r\n':
                    j += 1
                next_ch = text[j] if j < len(text) else ''
                if next_ch in ':,}]"' or j >= len(text):
                    # Closing quote
                    in_string = False
                    result.append(ch)
                else:
                    # Mid-string unescaped quote — replace with 「
                    result.append('「')
            i += 1
            continue
        result.append(ch)
        i += 1
    return ''.join(result)


def _parse_json(text: str) -> dict:
    """Find the top-level JSON object in an LLM response.

    Strategy order: strip fences → parse the top-level object → repair unescaped
    quotes and re-parse → lenient last-resort scan. Raises nothing — returns {}
    on total failure so the caller can decide how to surface the error.

    Strategies 1 and 2 decode at the FIRST ``{`` only (tolerating surrounding
    prose but NOT falling through to inner objects), so a malformed top-level
    object proceeds to the repair step instead of silently returning an inner
    object (e.g. a single ``core_claims`` item). The lenient first-decodable-
    object scan is kept only as a final fallback.
    """
    decoder = json.JSONDecoder()

    def _decode_at_first_brace(src: str) -> dict | None:
        """Decode the object at the first ``{`` only — no inner fall-through."""
        idx = src.find("{")
        if idx < 0:
            return None
        try:
            obj, _ = decoder.raw_decode(src, idx)
            return obj if isinstance(obj, dict) else None
        except Exception:  # noqa: BLE001
            return None

    def _first_obj(src: str) -> dict | None:
        """Lenient last resort: the first decodable object anywhere in src."""
        for i, ch in enumerate(src):
            if ch == "{":
                try:
                    obj, _ = decoder.raw_decode(src, i)
                    if isinstance(obj, dict):
                        return obj
                except Exception:  # noqa: BLE001
                    pass
        return None

    # Strategy 1: strip fences, parse the top-level object directly.
    clean = _strip_fences(text)
    result = _decode_at_first_brace(clean)
    if result is not None:
        return result

    # Strategy 2: repair unescaped quotes, then parse the top-level object.
    result = _decode_at_first_brace(_repair_json_quotes(clean))
    if result is not None:
        return result

    # Strategy 3 (last resort): lenient scan for the first decodable object.
    result = _first_obj(text)
    if result is not None:
        return result

    return {}


# ── Client + theme list helpers ──────────────────────────────────────────────
def get_llm_client():
    """Construct an Anthropic client from Streamlit secrets or env (fail-loud).

    Mirrors ``lib.llm_orchestrator._get_client`` so the page does not need to
    duplicate key resolution.
    """
    # The thesis feature is import-isolated, so unlike the Cockpit path it does
    # not transitively trigger another module's load_dotenv(). Load it explicitly
    # here (idempotent — only sets vars not already present) so a direct visit to
    # the Thesis Library page still picks up ANTHROPIC_API_KEY from .env. Explicit
    # repo-root path (CWD-independent), consistent with the other modules.
    load_dotenv(pathlib.Path(__file__).resolve().parent.parent.parent / ".env")

    import anthropic

    api_key = None
    try:
        import streamlit as st

        api_key = st.secrets.get("ANTHROPIC_API_KEY")
    except Exception:  # noqa: BLE001
        pass
    if not api_key:
        api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise ValueError("ANTHROPIC_API_KEY not configured")
    return anthropic.Anthropic(api_key=api_key)


def get_theme_names() -> list[str]:
    """Read the canonical theme-basket names (keys) read-only.

    Lazy import so loading this module never pulls heavy deps, and so the
    thesis feature degrades gracefully (empty list) if theme_baskets is absent.
    """
    try:
        try:
            from lib.theme_baskets import THEME_BASKETS  # type: ignore
        except Exception:  # noqa: BLE001
            from theme_baskets import THEME_BASKETS  # type: ignore
        return list(THEME_BASKETS.keys())
    except Exception as exc:  # noqa: BLE001
        _log.warning("get_theme_names: theme_baskets unavailable (%s)", exc)
        return []


# ── Document reading ─────────────────────────────────────────────────────────
def _decode_raw(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")


# Extensions we can extract text from, plus image extensions that are explicitly
# rejected with a friendly message at the UI layer (see pages/10) — never raised.
SUPPORTED_EXTS = ("txt", "md", "pdf", "docx", "pptx")
IMAGE_EXTS = ("jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp")


def is_image_ext(filename: str) -> bool:
    """True if *filename* has a recognised image extension."""
    return os.path.splitext(filename or "")[1].lower().lstrip(".") in IMAGE_EXTS


def read_document(file_bytes: bytes, filename: str) -> str:
    """Extract plain text from an uploaded document.

    Supported: ``.txt`` / ``.md`` (decoded directly), ``.pdf`` (pdfplumber),
    ``.docx`` (python-docx), ``.pptx`` (python-pptx) — all hard dependencies.
    Any other extension raises :class:`UnsupportedFormatError`. Image formats are
    intercepted at the UI layer with a friendly message (see ``is_image_ext``)
    before reaching here.
    """
    ext = os.path.splitext(filename or "")[1].lower().lstrip(".")

    if ext in ("txt", "md"):
        return _decode_raw(file_bytes)

    if ext == "pdf":
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            pages = [(p.extract_text() or "") for p in pdf.pages]
        return "\n".join(pages).strip() or _decode_raw(file_bytes)

    if ext == "docx":
        document = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in document.paragraphs).strip() or _decode_raw(file_bytes)

    if ext == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts).strip() or _decode_raw(file_bytes)

    raise UnsupportedFormatError(f"unsupported document format: {filename!r}")


# ── LLM call wrapper ─────────────────────────────────────────────────────────
def _call(llm_client, system: str, user: str, max_tokens: int) -> str:
    resp = llm_client.messages.create(
        model=_MODEL,
        max_tokens=max_tokens,
        system=system,
        messages=[{"role": "user", "content": user}],
    )
    return resp.content[0].text if (resp and resp.content) else ""


# ── 3a. Preview call ─────────────────────────────────────────────────────────
_PREVIEW_SYSTEM = """You are a research analyst assistant. Your task is to identify the distinct,
independently falsifiable arguments in the provided article or interview.

Output pure JSON with no markdown fences. The JSON must have exactly this structure:
{
  "arguments": [
    {
      "index": 1,
      "headline_en": "one sentence summarising this argument in English",
      "headline_zh": "Chinese one-sentence summary of the same argument using professional investment terminology",
      "primary_tickers": ["list of tickers explicitly named, uppercase, may be empty"],
      "primary_themes": ["main themes this argument relates to, may be empty"],
      "suggested_horizon": "short | mid | long"
    }
  ],
  "article_language": "zh | en | mixed",
  "article_title_en": "English title or best guess",
  "article_title_zh": "Chinese title or best guess using professional investment terminology"
}

An argument is independently falsifiable if it can be confirmed or refuted
without relying on other arguments in the same article.
Do not merge all arguments into one. Do not invent arguments not present in the text.
Output between 1 and 6 arguments. If the article has only one coherent argument, output 1.
Important: your response must be valid JSON. Use 「」 or 『』 for
quotation marks inside string values. Never use unescaped ASCII
double-quotes inside a JSON string value."""


def preview_article(file_text: str, llm_client) -> dict:
    """Lightweight call: list the distinct arguments the article makes.

    Returns the parsed dict (``arguments`` / ``article_language`` /
    ``article_title_*``). Raises :class:`ExtractionError` (with the raw response)
    on parse failure.
    """
    raw = _call(llm_client, _PREVIEW_SYSTEM, file_text, max_tokens=2000)
    parsed = _parse_json(raw)
    if not parsed or "arguments" not in parsed:
        raise ExtractionError(
            f"preview_article: could not parse arguments from response "
            f"({len(raw)} chars). Raw head: {raw[:400]!r}"
        )
    return parsed


# ── 3b. Full extraction call ─────────────────────────────────────────────────
# The source document may be in any language (Chinese, English, Japanese, etc.).
# Extract faithfully in the source language for verbatim fields (source_quote,
# assumptions, note). Output _en/_zh parallel keys for all bilingual prose fields
# regardless of source language. (Maintainer note — not sent to the LLM.)
_EXTRACT_SYSTEM_TMPL = """You are a research analyst assistant extracting a structured thesis card from a
research article or interview.

You are extracting argument #{argument_index}: "{argument_headline}"

Rules you must follow without exception:
- core_claims is the most important field. For the argument you are
  extracting, you MUST produce at least one item in core_claims. Each
  item represents one distinct claim the author makes within this
  argument. At minimum, produce one item summarising the argument's
  central thesis. Use claim_type "thesis" for the main claim,
  "risk_warning" for stated risks, "timing_call" for time-specific
  predictions, and "structural_observation" for factual observations
  the author uses as supporting evidence. claim_text_en and
  claim_text_zh are both required for every item.
- Your response must be valid JSON. When you need to quote a term,
  ticker, or phrase inside a string value, use Chinese quotation marks
  「」 or 『』 instead of ASCII double-quotes. Never use unescaped
  ASCII double-quotes inside a JSON string value.
- Faithfully represent what the author said. Do not add, infer, or embellish.
- For every numeric claim the author made explicitly, include it in numeric_claims
  with provenance "stated_by_author" and a source_quote.
- If the author implied a number without stating it, use provenance "inferred"
  and include source_quote showing what they said.
- If the author indicated a direction but gave no number (e.g. "margins will improve"),
  put it in unspecified_numerics. Never invent a proxy value.
  "direction" in unspecified_numerics must be exactly one of:
  "up", "down", or "unspecified". Never use other strings.
- For transmission_chain, each step must have a clear from_node and to_node.
  If you are inferring a step the author implied but did not state,
  set provenance to "inferred".
- For confirmation_conditions and falsification_conditions, if the author stated
  explicit conditions, set provenance "stated_by_author".
  If you are inferring them from the argument logic, set provenance "inferred".
  "observable" must be exactly one of these three strings:
  "machine_checkable", "human_judgment", or "unspecified".
  Never use a boolean. Use "machine_checkable" when the condition
  can be verified from market data or financial statements.
  Use "human_judgment" when it requires qualitative assessment.
  Use "unspecified" when unclear.
- "affected_horizons" must only contain values from this exact list:
  ["short", "mid", "long"]. Never use Chinese or descriptive strings.
  short = weeks to 3 months, mid = 3 months to 18 months,
  long = 18 months or more.
- current_evidence_status must always be "unknown". Do not change this.
- evidence_refs must always be an empty list. Do not add anything.
- coi.status must always be "coi_unassessed" unless the author explicitly
  discloses a conflict of interest in the text, in which case set it to
  "coi_disclosed" and quote the disclosure in coi.notes.
- Bilingual prose fields: output BOTH an English version (field_en) and a Chinese
  version (field_zh) for claim_text, mechanism, condition_text,
  event_or_hypothesis, and scenario notes. The Chinese must use professional
  investment research terminology, not machine translation.
- Verbatim author fields stay single-language — quote the author's original text
  and do NOT translate: source_quote (numeric_claims), note (unspecified_numerics),
  and each assumptions entry.
- affected_themes must only contain values from this list: {theme_list}
  Themes the author mentions that are not in this list go into unmapped_themes.

Output pure JSON with no markdown fences matching the ThesisCard schema exactly.
Use these field names: core_claims (each with claim_text_en, claim_text_zh,
claim_type, related_tickers, related_themes), numeric_claims (metric, value, unit,
applies_to, time_reference, provenance, source_quote), unspecified_numerics (metric,
direction, note), assumptions (list of strings), scenarios (each with
event_or_hypothesis_en, event_or_hypothesis_zh, transmission_chain [step, from_node,
to_node, mechanism_en, mechanism_zh, provenance], affected_horizons, affected_themes,
unmapped_themes, affected_tickers, confirmation_conditions [condition_text_en,
condition_text_zh, observable, provenance], falsification_conditions [same],
notes_en, notes_zh), and coi (status, notes)."""


def _normalise_str(s) -> str:
    return s if isinstance(s, str) else ""


def _scenario_event_text(scenario: dict) -> str:
    """Canonical event text for id/hashing — English preferred, then bare, then zh."""
    return (
        _normalise_str(scenario.get("event_or_hypothesis_en"))
        or _normalise_str(scenario.get("event_or_hypothesis"))
        or _normalise_str(scenario.get("event_or_hypothesis_zh"))
    )


def _scenario_id(scenario: dict) -> str:
    """sha256 (first 16 chars) of the scenario's normalised salient content."""
    chain = scenario.get("transmission_chain") or []
    chain_repr = "|".join(
        f"{(st.get('from_node') or '').strip()}->{(st.get('to_node') or '').strip()}"
        for st in chain if isinstance(st, dict)
    )
    basis = (
        _scenario_event_text(scenario).strip().lower()
        + "::" + chain_repr.lower()
    )
    return hashlib.sha256(basis.encode("utf-8")).hexdigest()[:16]


def _coi_from_llm(raw_coi) -> dict:
    """Accept a disclosed COI from the model; otherwise force unassessed."""
    if isinstance(raw_coi, dict) and raw_coi.get("status") == "coi_disclosed":
        return {"status": "coi_disclosed", "notes": _normalise_str(raw_coi.get("notes"))}
    return empty_coi()


def _normalise_observable(val) -> str:
    if val is True or val == "true":
        return "machine_checkable"
    if val is False or val == "false":
        return "unspecified"
    if val in {"machine_checkable", "human_judgment", "unspecified"}:
        return val
    return "unspecified"


def _normalise_horizon(val: str) -> str | None:
    v = str(val).lower().strip()
    if v in {"short", "short-term", "短线", "短期"}:
        return "short"
    if v in {"mid", "medium", "mid-term", "medium-term", "中线", "中期"}:
        return "mid"
    if v in {"long", "long-term", "长线", "长期"}:
        return "long"
    return None  # drop unrecognised values


def _normalise_direction(val) -> str:
    v = str(val).lower().strip()
    if v in {"up", "upside", "upside_risk", "increase", "higher",
             "上升", "上行", "上涨"}:
        return "up"
    if v in {"down", "downside", "downside_risk", "decrease", "lower",
             "下降", "下行", "下跌"}:
        return "down"
    return "unspecified"


def _normalise_scenarios(raw_scenarios, card_themes_fallback=None) -> list[dict]:
    """Force scenario invariants (schema_version, evidence fields, id)."""
    out: list[dict] = []
    if not isinstance(raw_scenarios, list):
        return out
    for sc in raw_scenarios:
        if not isinstance(sc, dict):
            continue
        norm = dict(sc)
        norm["schema_version"] = SCHEMA_VERSION_SCENARIO
        # Invariants — never trusted from the model. Overwrite, but warn first so
        # a misbehaving LLM is observable rather than silently corrected.
        if norm.get("current_evidence_status") != "unknown":
            _log.warning(
                "LLM returned current_evidence_status=%r on scenario %r — "
                "overriding to 'unknown'",
                norm.get("current_evidence_status"),
                norm.get("scenario_id", "?"),
            )
        norm["current_evidence_status"] = "unknown"
        ev = norm.get("evidence_refs")
        if ev != []:
            _log.warning(
                "LLM returned non-empty or non-list evidence_refs on scenario %r — "
                "clearing to []",
                norm.get("scenario_id", "?"),
            )
        norm["evidence_refs"] = []
        # Structural defaults for list fields.
        for key in (
            "transmission_chain", "affected_horizons", "affected_themes",
            "unmapped_themes", "affected_tickers",
            "confirmation_conditions", "falsification_conditions",
        ):
            if not isinstance(norm.get(key), list):
                norm[key] = []
        # Coerce affected_horizons to the {short, mid, long} enum; drop unknowns.
        norm["affected_horizons"] = [
            h for raw in norm.get("affected_horizons", [])
            if (h := _normalise_horizon(raw)) is not None
        ]
        # Coerce each condition's observable to the enum (e.g. boolean true → string).
        for cond_field in ("confirmation_conditions", "falsification_conditions"):
            for cond in norm.get(cond_field, []):
                if isinstance(cond, dict):
                    cond["observable"] = _normalise_observable(cond.get("observable"))
        # Normalise bilingual scenario prose (en/zh); keep any legacy bare key.
        for base in ("event_or_hypothesis", "notes"):
            for key in (f"{base}_en", f"{base}_zh", base):
                if key in norm:
                    norm[key] = _normalise_str(norm.get(key))
        norm["scenario_id"] = _scenario_id(norm)
        # Uppercase any tickers.
        norm["affected_tickers"] = [
            str(tk).upper() for tk in norm.get("affected_tickers", []) if str(tk).strip()
        ]
        out.append(norm)
    return out


def extract_card(
    file_text: str,
    argument_index: int,
    argument_headline: str,
    horizon_type: str,
    doc_meta: dict,
    llm_client,
    extraction_seq: int = 1,
) -> dict:
    """Extract one structured thesis card for a chosen argument.

    Returns an assembled (but **unvalidated**) ThesisCard dict. The LLM supplies
    only the interpretive content (claims, numeric_claims, scenarios, ...); every
    structural / invariant field is set deterministically here. Raises
    :class:`ExtractionError` if the model response is not parseable JSON.

    ``doc_meta`` must contain: doc_hash, doc_path, doc_type, author,
    author_affiliation, publication_date, publication_date_provenance, language,
    title.
    """
    theme_list = ", ".join(get_theme_names()) or "(no theme baskets configured)"
    system = _EXTRACT_SYSTEM_TMPL.format(
        argument_index=argument_index,
        argument_headline=argument_headline,
        theme_list=theme_list,
    )
    raw_text = _call(llm_client, system, file_text, max_tokens=4000)
    parsed = _parse_json(raw_text)

    # Guard against mis-parse: if the scanner fell through to an inner
    # object, parsed will look like a claim item rather than a card.
    # A valid top-level response must contain at least one of these keys.
    _EXPECTED_TOP_LEVEL_KEYS = {
        "core_claims", "numeric_claims", "scenarios",
        "assumptions", "coi", "source_language"
    }
    if parsed and not (set(parsed.keys()) & _EXPECTED_TOP_LEVEL_KEYS):
        # The parse landed on an inner object. Log the raw response and raise.
        _log.error(
            "extract_card: _parse_json returned an inner object (keys=%s). "
            "Raw response snippet: %.500s",
            list(parsed.keys()),
            raw_text[:500] if raw_text else "",
        )
        raise ExtractionError(
            f"JSON parse landed on inner object (keys={list(parsed.keys())}). "
            f"The LLM likely emitted unescaped quotes. Raw (first 500 chars): "
            f"{raw_text[:500] if raw_text else ''}"
        )

    if not parsed:
        raise ExtractionError(
            f"Failed to parse LLM response as JSON. "
            f"Raw (first 500 chars): {raw_text[:500] if raw_text else ''}"
        )

    doc_hash = str(doc_meta.get("doc_hash", ""))
    card_id = card_id_from_hash(doc_hash, extraction_seq)
    now_iso = datetime.now().isoformat(timespec="seconds")

    # ── core_claims: deterministic claim_id; pass interpretive text through ──
    core_claims: list[dict] = []
    for n, cc in enumerate(parsed.get("core_claims") or [], start=1):
        if not isinstance(cc, dict):
            continue
        core_claims.append({
            "claim_id": f"{card_id}-c{n}",
            "claim_text_en": _normalise_str(cc.get("claim_text_en")),
            "claim_text_zh": _normalise_str(cc.get("claim_text_zh")),
            "claim_type": _normalise_str(cc.get("claim_type")) or "thesis",
            "related_tickers": [
                str(tk).upper() for tk in (cc.get("related_tickers") or [])
                if str(tk).strip()
            ],
            "related_themes": list(cc.get("related_themes") or []),
        })

    # ── numeric_claims / unspecified_numerics: content as the author stated ──
    numeric_claims = [nc for nc in (parsed.get("numeric_claims") or []) if isinstance(nc, dict)]
    unspecified = [un for un in (parsed.get("unspecified_numerics") or []) if isinstance(un, dict)]
    # Coerce each direction to the {up, down, unspecified} enum.
    for un in unspecified:
        un["direction"] = _normalise_direction(un.get("direction"))

    assumptions = [str(a) for a in (parsed.get("assumptions") or []) if str(a).strip()]
    scenarios = _normalise_scenarios(parsed.get("scenarios"))

    card: dict = {
        "schema_version": SCHEMA_VERSION_THESIS,
        "card_id": card_id,
        "source": {
            "doc_hash": doc_hash,
            "doc_path": str(doc_meta.get("doc_path", "")),
            "doc_type": str(doc_meta.get("doc_type", "other")),
            "title": str(doc_meta.get("title", "")),
            "author": str(doc_meta.get("author", "unknown")) or "unknown",
            "author_affiliation": str(doc_meta.get("author_affiliation", "unknown")) or "unknown",
            "publication_date": doc_meta.get("publication_date"),
            "publication_date_provenance": str(
                doc_meta.get("publication_date_provenance", "unspecified")
            ),
            "language": str(doc_meta.get("language", "en")),
            "ingested_at": now_iso,
        },
        "horizon_type": horizon_type,
        "coi": _coi_from_llm(parsed.get("coi")),
        "card_status": "active",
        "core_claims": core_claims,
        "numeric_claims": numeric_claims,
        "unspecified_numerics": unspecified,
        "assumptions": assumptions,
        "scenarios": scenarios,
        "extraction_meta": {
            "llm_model": _MODEL,
            "prompt_version": PROMPT_VERSION,
            "extracted_at": now_iso,
            "extraction_seq": int(extraction_seq),
        },
    }
    return card


__all__ = [
    "ExtractionError",
    "UnsupportedFormatError",
    "preview_article",
    "extract_card",
    "read_document",
    "get_llm_client",
    "get_theme_names",
]
